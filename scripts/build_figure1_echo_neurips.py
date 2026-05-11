# -*- coding: utf-8 -*-
from __future__ import annotations

import shutil
from pathlib import Path
from zipfile import ZipFile

import fitz
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PAPER_PDF = Path(r"E:\edge\VLA_Memory_For_NeurIPS_2026 (2).pdf")
OUT_DIR = ROOT / "thesis_assets_generated"
ASSET_DIR = OUT_DIR / "figure1_echo_neurips_assets"
EXPORT_DIR = ROOT / "debug" / "figure1_echo_neurips_export"

HYPERBOLIC_SRC = OUT_DIR / "figure1_echo_hyperbolic_image2_neurips.png"
HYPERBOLIC_CIRCLE = ASSET_DIR / "hyperbolic_image2_neurips_circle.png"

OUT_PPTX = OUT_DIR / "figure1_echo_neurips.pptx"
OUT_PDF = OUT_DIR / "figure1_echo_neurips.pdf"
OUT_PNG = OUT_DIR / "figure1_echo_neurips.png"

SLIDE_W = 13.2
SLIDE_H = 4.95

COLORS = {
    "black": (30, 30, 30),
    "muted": (95, 95, 95),
    "line": (190, 190, 190),
    "line_dark": (82, 82, 82),
    "blue": (55, 105, 170),
    "blue_soft": (232, 241, 252),
    "amber": (201, 128, 25),
    "amber_soft": (255, 244, 225),
    "green": (46, 145, 74),
    "green_soft": (232, 248, 237),
    "red": (185, 60, 60),
    "red_soft": (252, 236, 236),
    "gray_soft": (247, 247, 247),
    "white": (255, 255, 255),
}


def rgb(name: str) -> RGBColor:
    return RGBColor(*COLORS[name])


def inch(value: float):
    return Inches(value)


def reset_dir(path: Path) -> None:
    resolved = path.resolve()
    allowed = (ROOT / "debug").resolve()
    if not resolved.is_relative_to(allowed):
        raise RuntimeError(f"Refusing to clear unexpected directory: {resolved}")
    shutil.rmtree(resolved, ignore_errors=True)
    resolved.mkdir(parents=True, exist_ok=True)


def safe_unlink(path: Path) -> None:
    try:
        if path.exists():
            path.unlink()
    except PermissionError:
        alt = path.with_name(path.stem + "_new" + path.suffix)
        if alt.exists():
            alt.unlink()
        globals()[f"OUT_{path.suffix.upper()[1:]}"] = alt


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float = 7.2,
    color: str = "black",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    box.shadow.inherit = False
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.02)
    tf.margin_right = inch(0.02)
    tf.margin_top = inch(0.00)
    tf.margin_bottom = inch(0.00)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str = "white",
    line: str | None = "line",
    width: float = 0.55,
):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(width)
    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str = "line_dark",
    width: float = 0.55,
    arrow: bool = False,
    dash: bool = False,
):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    line.shadow.inherit = False
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if arrow:
        line.line.end_arrowhead = 3
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def add_node(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    fill: str,
    line: str,
    size: float = 6.4,
    bold: bool = False,
):
    add_rect(slide, x, y, w, h, fill=fill, line=line, width=0.6)
    add_text(slide, x + 0.02, y + 0.01, w - 0.04, h - 0.02, label, size=size, bold=bold, align=PP_ALIGN.CENTER)


def extract_banana_frames() -> list[Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(PAPER_PDF))
    images = doc[12].get_images(full=True)
    frames: list[Path] = []
    for idx, image_no in enumerate([5, 6, 7, 8], start=1):
        pix = fitz.Pixmap(doc, images[image_no - 1][0])
        if pix.n - pix.alpha > 3:
            pix = fitz.Pixmap(fitz.csRGB, pix)
        raw = ASSET_DIR / f"banana_raw_{idx}.png"
        pix.save(str(raw))
        im = Image.open(raw).convert("RGB")
        w, h = im.size
        crop = im.crop((0, int(h * 0.157), w, int(h * 0.754)))
        cw, ch = crop.size
        target_ratio = 4 / 3
        if cw / ch > target_ratio:
            nw = int(ch * target_ratio)
            left = (cw - nw) // 2
            crop = crop.crop((left, 0, left + nw, ch))
        else:
            nh = int(cw / target_ratio)
            top = (ch - nh) // 2
            crop = crop.crop((0, top, cw, top + nh))
        out = ASSET_DIR / f"banana_frame_{idx}.png"
        crop.resize((960, 720), Image.Resampling.LANCZOS).save(out, quality=95)
        frames.append(out)
    return frames


def prepare_hyperbolic_asset() -> Path:
    if not HYPERBOLIC_SRC.exists():
        raise FileNotFoundError(f"Missing generated hyperbolic image: {HYPERBOLIC_SRC}")
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    im = Image.open(HYPERBOLIC_SRC).convert("RGBA")
    w, h = im.size
    pad = int(min(w, h) * 0.025)
    mask = Image.new("L", (w, h), 0)
    draw = ImageDraw.Draw(mask)
    draw.ellipse((pad, pad, w - pad, h - pad), fill=255)
    im.putalpha(mask)
    im.save(HYPERBOLIC_CIRCLE)
    return HYPERBOLIC_CIRCLE


def build_left(slide, frames: list[Path]) -> None:
    add_text(slide, 0.25, 0.12, 5.95, 0.20, "(a) Embodied task hierarchy from real executions", 8.2, bold=True)
    add_rect(slide, 0.22, 0.39, 6.12, 4.30, fill="white", line="line")

    img_w, img_h, gap = 1.36, 1.02, 0.08
    for i, frame in enumerate(frames):
        x = 0.34 + i * (img_w + gap)
        pic = slide.shapes.add_picture(str(frame), inch(x), inch(0.55), width=inch(img_w), height=inch(img_h))
        pic.shadow.inherit = False
        add_rect(slide, x, 0.55, img_w, img_h, fill="white", line="line", width=0.45).fill.background()
        add_text(slide, x + 0.03, 0.46, 0.35, 0.12, f"t{i + 1}", 5.8, "muted", bold=True, align=PP_ALIGN.LEFT)
        if i < 3:
            add_line(slide, x + img_w + 0.015, 1.06, x + img_w + gap - 0.015, 1.06, "line_dark", 0.5, arrow=True)

    add_text(slide, 0.55, 1.70, 5.45, 0.18, "Real-world frames expose a macro-to-micro containment structure.", 6.5, "muted", align=PP_ALIGN.CENTER)

    add_line(slide, 0.70, 2.05, 0.70, 4.27, "line_dark", 0.55, arrow=True)
    add_text(slide, 0.42, 2.02, 0.55, 0.14, "macro", 5.8, "muted", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 0.42, 4.31, 0.55, 0.14, "micro", 5.8, "muted", bold=True, align=PP_ALIGN.CENTER)

    add_node(slide, 2.50, 2.05, 1.65, 0.26, "Long-horizon manipulation", "green_soft", "green", 6.1, True)
    add_node(slide, 1.30, 2.73, 1.22, 0.25, "Pick banana", "blue_soft", "blue", 6.1, True)
    add_node(slide, 4.05, 2.73, 1.22, 0.25, "Place in bowl", "amber_soft", "amber", 6.1, True)
    add_line(slide, 3.32, 2.31, 1.91, 2.73, "line_dark", 0.5)
    add_line(slide, 3.32, 2.31, 4.66, 2.73, "line_dark", 0.5)

    leaves = [
        (1.03, "Approach", "blue_soft", "blue"),
        (1.78, "Grasp", "blue_soft", "blue"),
        (2.53, "Lift", "blue_soft", "blue"),
        (3.78, "Transport", "amber_soft", "amber"),
        (4.63, "Align", "amber_soft", "amber"),
        (5.38, "Release", "amber_soft", "amber"),
    ]
    for x, label, fill, line in leaves:
        add_node(slide, x, 3.62, 0.62, 0.24, label, fill, line, 5.5)
    for x in [1.34, 2.09, 2.84]:
        add_line(slide, 1.91, 2.98, x, 3.62, "line_dark", 0.45)
    for x in [4.09, 4.94, 5.69]:
        add_line(slide, 4.66, 2.98, x, 3.62, "line_dark", 0.45)

    add_text(slide, 1.20, 4.17, 4.85, 0.22, "Fine-grained controls are reusable leaves under semantic subgoals.", 6.4, "black", align=PP_ALIGN.CENTER)


def build_right(slide, hyperbolic: Path) -> None:
    add_text(slide, 6.65, 0.12, 6.20, 0.20, "(b) Flat memory versus ECHO hyperbolic memory", 8.2, bold=True)
    add_rect(slide, 6.60, 0.39, 6.38, 1.72, fill="white", line="line")
    add_rect(slide, 6.83, 0.61, 0.78, 0.20, fill="red_soft", line="red", width=0.55)
    add_text(slide, 6.88, 0.64, 0.68, 0.10, "Baseline", 5.8, "red", True, PP_ALIGN.CENTER)
    add_text(slide, 7.78, 0.61, 1.65, 0.20, "Linear flat storage", 7.6, bold=True)

    qx, qy = 6.96, 1.08
    for i, label in enumerate(["f1", "f2", "f3", "f4", "f5", "...", "fn"]):
        fill = "red_soft" if label in {"f2", "..."} else "gray_soft"
        add_rect(slide, qx + i * 0.52, qy, 0.34, 0.24, fill=fill, line="line", width=0.45)
        add_text(slide, qx + i * 0.52, qy + 0.03, 0.34, 0.11, label, 5.8, align=PP_ALIGN.CENTER)
        if i < 6:
            add_line(slide, qx + i * 0.52 + 0.34, qy + 0.12, qx + (i + 1) * 0.52, qy + 0.12, "line_dark", 0.35)
    add_line(slide, 7.07, 1.55, 10.18, 1.55, "red", 0.55, arrow=True, dash=True)
    add_text(slide, 7.45, 1.64, 1.04, 0.14, "global scan", 5.8, "red", True, PP_ALIGN.CENTER)
    add_text(slide, 9.45, 1.64, 0.55, 0.14, "O(N)", 5.8, "red", True, PP_ALIGN.CENTER)
    add_text(slide, 10.42, 1.42, 1.16, 0.20, "weak hierarchy", 6.0, "red", True)

    add_rect(slide, 6.60, 2.29, 6.38, 2.40, fill="white", line="line")
    add_rect(slide, 6.83, 2.50, 0.54, 0.20, fill="green_soft", line="green", width=0.55)
    add_text(slide, 6.88, 2.53, 0.44, 0.10, "ECHO", 5.8, "green", True, PP_ALIGN.CENTER)
    add_text(slide, 7.55, 2.50, 2.18, 0.20, "Hyperbolic continuous space", 7.6, bold=True)

    pic = slide.shapes.add_picture(str(hyperbolic), inch(6.88), inch(2.82), width=inch(1.62), height=inch(1.62))
    pic.shadow.inherit = False
    add_text(slide, 7.07, 4.45, 1.24, 0.14, "Poincare disk", 5.8, "muted", align=PP_ALIGN.CENTER)

    add_node(slide, 9.12, 3.00, 1.38, 0.24, "semantic tree", "green_soft", "green", 5.8, True)
    add_node(slide, 10.95, 3.00, 1.36, 0.24, "top-down retrieval", "blue_soft", "blue", 5.8, True)
    add_line(slide, 8.52, 3.45, 9.12, 3.12, "line_dark", 0.5, arrow=True)
    add_line(slide, 10.50, 3.12, 10.95, 3.12, "line_dark", 0.5, arrow=True)

    add_rect(slide, 9.14, 3.63, 2.90, 0.45, fill="gray_soft", line="line", width=0.45)
    add_text(slide, 9.23, 3.70, 2.72, 0.13, "Hyperbolic entailment cones", 6.2, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 9.23, 3.90, 2.72, 0.12, "partial containment for overlapping actions", 5.6, "muted", align=PP_ALIGN.CENTER)
    add_line(slide, 8.50, 3.66, 9.14, 3.84, "line_dark", 0.5, arrow=True)


def build_presentation(frames: list[Path], hyperbolic: Path) -> None:
    prs = Presentation()
    prs.slide_width = inch(SLIDE_W)
    prs.slide_height = inch(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill="white", line=None)
    build_left(slide, frames)
    build_right(slide, hyperbolic)
    prs.save(OUT_PPTX)


def export_powerpoint() -> None:
    import win32com.client

    reset_dir(EXPORT_DIR)
    for path in [OUT_PDF, OUT_PNG]:
        if path.exists():
            path.unlink()

    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = True
    app.DisplayAlerts = 0
    presentation = app.Presentations.Open(str(OUT_PPTX), False, True, False)
    try:
        presentation.SaveAs(str(OUT_PDF), 32)
        presentation.Export(str(EXPORT_DIR), "PNG", 3600, 1350)
    finally:
        presentation.Close()
        app.Quit()

    exported = sorted(EXPORT_DIR.glob("*.PNG")) + sorted(EXPORT_DIR.glob("*.png"))
    if exported:
        shutil.copy2(exported[0], OUT_PNG)


def verify_ppt_text() -> None:
    texts: list[str] = []
    with ZipFile(OUT_PPTX) as z:
        for name in z.namelist():
            if name.startswith("ppt/slides/") and name.endswith(".xml"):
                xml = z.read(name).decode("utf-8", errors="ignore")
                parts = xml.split("<a:t>")[1:]
                texts.extend(part.split("</a:t>")[0] for part in parts)
    joined = " ".join(texts)
    non_ascii = [t for t in texts if not t.isascii()]
    if non_ascii:
        raise RuntimeError(f"Non-English/non-ASCII text found: {non_ascii}")
    if "fruit salad" in joined.lower():
        raise RuntimeError("Forbidden fruit-salad text found.")


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    for path in [OUT_PPTX, OUT_PDF, OUT_PNG]:
        if path.exists():
            path.unlink()
    frames = extract_banana_frames()
    hyperbolic = prepare_hyperbolic_asset()
    build_presentation(frames, hyperbolic)
    verify_ppt_text()
    export_powerpoint()
    print(f"PPTX: {OUT_PPTX}")
    print(f"PDF:  {OUT_PDF}")
    print(f"PNG:  {OUT_PNG}")
    print(f"Hyperbolic asset: {HYPERBOLIC_SRC}")


if __name__ == "__main__":
    main()
