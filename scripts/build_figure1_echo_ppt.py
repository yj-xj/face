# -*- coding: utf-8 -*-
from __future__ import annotations

import math
import shutil
from pathlib import Path

import fitz
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PAPER_PDF = Path(r"E:\edge\VLA_Memory_For_NeurIPS_2026 (2).pdf")
ASSET_DIR = ROOT / "thesis_assets_generated" / "figure1_echo_assets"
OUT_DIR = ROOT / "thesis_assets_generated"
OUT_PPTX = OUT_DIR / "figure1_echo_memory_image2.pptx"
OUT_PDF = OUT_DIR / "figure1_echo_memory_image2.pdf"
OUT_PNG = OUT_DIR / "figure1_echo_memory_image2.png"
EXPORT_DIR = ROOT / "debug" / "figure1_echo_export"
HYPERBOLIC_IMAGE2 = OUT_DIR / "figure1_echo_hyperbolic_image2.png"
HYPERBOLIC_IMAGE2_CIRCLE = OUT_DIR / "figure1_echo_hyperbolic_image2_circle.png"

SLIDE_W = 13.333
SLIDE_H = 5.20

COLORS = {
    "ink": (31, 41, 55),
    "muted": (92, 105, 119),
    "line": (203, 213, 225),
    "line_dark": (116, 133, 150),
    "panel": (250, 252, 254),
    "panel_warm": (255, 251, 245),
    "panel_green": (246, 253, 249),
    "blue": (37, 99, 235),
    "blue_soft": (221, 235, 255),
    "amber": (245, 158, 11),
    "amber_soft": (255, 239, 204),
    "green": (22, 163, 74),
    "green_soft": (220, 252, 231),
    "red": (197, 48, 48),
    "red_soft": (254, 226, 226),
    "gray_soft": (241, 245, 249),
    "white": (255, 255, 255),
}


def c(name: str) -> RGBColor:
    return RGBColor(*COLORS[name])


def inch(value: float):
    return Inches(value)


def clear_export_dir(path: Path) -> None:
    resolved = path.resolve()
    allowed = (ROOT / "debug").resolve()
    if not resolved.is_relative_to(allowed):
        raise RuntimeError(f"Refusing to clear unexpected directory: {resolved}")
    shutil.rmtree(resolved, ignore_errors=True)
    resolved.mkdir(parents=True, exist_ok=True)


def safe_unlink(path: Path) -> None:
    if path.exists():
        path.unlink()


def extract_banana_frames() -> list[Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(PAPER_PDF))
    page = doc[12]
    images = page.get_images(full=True)
    if len(images) < 8:
        raise RuntimeError("Expected at least eight real-world example images on page 13.")

    frame_paths: list[Path] = []
    for idx, image_no in enumerate([5, 6, 7, 8], start=1):
        xref = images[image_no - 1][0]
        pix = fitz.Pixmap(doc, xref)
        if pix.n - pix.alpha > 3:
            pix = fitz.Pixmap(fitz.csRGB, pix)
        raw_path = ASSET_DIR / f"banana_raw_{idx}.png"
        pix.save(str(raw_path))

        im = Image.open(raw_path).convert("RGB")
        w, h = im.size
        crop = im.crop((0, int(h * 0.157), w, int(h * 0.754)))
        cw, ch = crop.size
        target_ratio = 4 / 3
        if cw / ch > target_ratio:
            new_w = int(ch * target_ratio)
            left = (cw - new_w) // 2
            crop = crop.crop((left, 0, left + new_w, ch))
        else:
            new_h = int(cw / target_ratio)
            top = (ch - new_h) // 2
            crop = crop.crop((0, top, cw, top + new_h))
        crop = crop.resize((960, 720), Image.Resampling.LANCZOS)
        out_path = ASSET_DIR / f"banana_frame_{idx}.png"
        crop.save(out_path, quality=95)
        frame_paths.append(out_path)
    return frame_paths


def prepare_hyperbolic_image2() -> Path | None:
    if not HYPERBOLIC_IMAGE2.exists():
        return None
    im = Image.open(HYPERBOLIC_IMAGE2).convert("RGBA")
    w, h = im.size
    pad = int(min(w, h) * 0.026)
    mask = Image.new("L", (w, h), 0)
    draw = ImageDraw.Draw(mask)
    draw.ellipse((pad, pad, w - pad, h - pad), fill=255)
    im.putalpha(mask)
    im.save(HYPERBOLIC_IMAGE2_CIRCLE)
    return HYPERBOLIC_IMAGE2_CIRCLE


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float = 8.5,
    color: str = "ink",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font_name: str = "Arial",
):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.04)
    tf.margin_right = inch(0.04)
    tf.margin_top = inch(0.01)
    tf.margin_bottom = inch(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = c(color)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str = "white",
    line: str | None = "line",
    radius: bool = False,
    transparency: int = 0,
):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c(fill)
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = c(line)
        shape.line.width = Pt(0.7)
    return shape


def add_node(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    fill: str,
    line: str,
    size: float = 7.5,
    bold: bool = False,
):
    add_rect(slide, x, y, w, h, fill=fill, line=line, radius=True)
    return add_text(slide, x + 0.02, y + 0.01, w - 0.04, h - 0.02, text, size, "ink", bold, PP_ALIGN.CENTER)


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str = "line_dark",
    width: float = 0.8,
    arrow: bool = False,
    dash: bool = False,
):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    line.line.color.rgb = c(color)
    line.line.width = Pt(width)
    if arrow:
        line.line.end_arrowhead = 3
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def add_picture_with_border(slide, path: Path, x: float, y: float, w: float, h: float):
    slide.shapes.add_picture(str(path), inch(x), inch(y), width=inch(w), height=inch(h))
    border = add_rect(slide, x, y, w, h, fill="white", line="line", radius=False, transparency=100)
    border.fill.background()
    return border


def add_cross(slide, x: float, y: float, size: float = 0.13, color: str = "red"):
    add_line(slide, x, y, x + size, y + size, color=color, width=1.7)
    add_line(slide, x + size, y, x, y + size, color=color, width=1.7)


def add_freeform_polygon(slide, pts: list[tuple[float, float]], fill: str, line: str | None, transparency: int):
    start_x, start_y = pts[0]
    fb = slide.shapes.build_freeform(inch(start_x), inch(start_y))
    fb.add_line_segments([(inch(x), inch(y)) for x, y in pts[1:]], close=True)
    shape = fb.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = c(fill)
    shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = c(line)
        shape.line.width = Pt(0.6)
    return shape


def disk_point(cx: float, cy: float, r: float, px: float, py: float) -> tuple[float, float]:
    return cx + px * r, cy + py * r


def boundary_arc(cx: float, cy: float, r: float, start_deg: float, end_deg: float, steps: int = 18):
    pts = []
    for i in range(steps + 1):
        t = math.radians(start_deg + (end_deg - start_deg) * i / steps)
        pts.append((cx + math.cos(t) * r, cy + math.sin(t) * r))
    return pts


def add_poincare_disk(slide, cx: float, cy: float, r: float):
    disk_bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(cx - r), inch(cy - r), inch(2 * r), inch(2 * r))
    disk_bg.fill.solid()
    disk_bg.fill.fore_color.rgb = c("white")
    disk_bg.line.fill.background()

    # Entailment cones sit above the disk background and below the metric grid.
    pick_node = disk_point(cx, cy, r, -0.38, 0.12)
    place_node = disk_point(cx, cy, r, 0.38, 0.12)
    pick_poly = [pick_node] + boundary_arc(cx, cy, r * 0.93, 138, 286)
    place_poly = [place_node] + boundary_arc(cx, cy, r * 0.93, 254, 402)
    add_freeform_polygon(slide, pick_poly, fill="blue", line=None, transparency=60)
    add_freeform_polygon(slide, place_poly, fill="amber", line=None, transparency=60)

    for rr in [0.28, 0.48, 0.68, 0.86]:
        ring = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            inch(cx - r * rr),
            inch(cy - r * rr),
            inch(2 * r * rr),
            inch(2 * r * rr),
        )
        ring.fill.background()
        ring.line.color.rgb = c("line")
        ring.line.width = Pt(0.4)

    for deg in range(0, 360, 30):
        t = math.radians(deg)
        add_line(
            slide,
            cx + math.cos(t) * r * 0.12,
            cy + math.sin(t) * r * 0.12,
            cx + math.cos(t) * r * 0.90,
            cy + math.sin(t) * r * 0.90,
            color="line",
            width=0.35,
        )

    points = {
        "root": (0.00, -0.02, "green", 0.070),
        "pick": (-0.38, 0.12, "blue", 0.065),
        "place": (0.38, 0.12, "amber", 0.065),
        "p1": (-0.72, 0.44, "blue", 0.045),
        "p2": (-0.47, 0.63, "blue", 0.045),
        "p3": (-0.18, 0.58, "green", 0.052),
        "q1": (0.18, 0.58, "green", 0.052),
        "q2": (0.47, 0.63, "amber", 0.045),
        "q3": (0.72, 0.44, "amber", 0.045),
    }

    for a, b in [("root", "pick"), ("root", "place"), ("pick", "p1"), ("pick", "p2"), ("pick", "p3"), ("place", "q1"), ("place", "q2"), ("place", "q3")]:
        ax, ay, _, _ = points[a]
        bx, by, _, _ = points[b]
        x1, y1 = disk_point(cx, cy, r, ax, ay)
        x2, y2 = disk_point(cx, cy, r, bx, by)
        add_line(slide, x1, y1, x2, y2, color="line_dark", width=0.55)

    for _, (px, py, color, size) in points.items():
        x, y = disk_point(cx, cy, r, px, py)
        node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(x - size), inch(y - size), inch(2 * size), inch(2 * size))
        node.fill.solid()
        node.fill.fore_color.rgb = c(color)
        node.line.color.rgb = c("white")
        node.line.width = Pt(0.9)

    boundary = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, inch(cx - r), inch(cy - r), inch(2 * r), inch(2 * r))
    boundary.fill.background()
    boundary.line.color.rgb = c("line_dark")
    boundary.line.width = Pt(1.0)


def build_left_panel(slide, frames: list[Path]):
    add_rect(slide, 0.18, 0.18, 6.18, 4.84, fill="panel", line="line")
    add_text(slide, 0.30, 0.23, 5.80, 0.24, "(a) Natural hierarchy in embodied manipulation", 9.5, "ink", True)

    img_x = 0.34
    img_y = 0.54
    img_w = 1.37
    img_h = 1.03
    gap = 0.09
    for i, frame in enumerate(frames):
        x = img_x + i * (img_w + gap)
        add_picture_with_border(slide, frame, x, img_y, img_w, img_h)
        add_rect(slide, x + 0.06, img_y + 0.06, 0.24, 0.15, fill="white", line=None, transparency=10)
        add_text(slide, x + 0.06, img_y + 0.055, 0.24, 0.15, f"t{i + 1}", 6.2, "muted", True, PP_ALIGN.CENTER)
        if i < 3:
            add_line(slide, x + img_w + 0.02, img_y + img_h / 2, x + img_w + gap - 0.02, img_y + img_h / 2, color="muted", width=0.9, arrow=True)

    add_text(slide, 0.53, 1.66, 5.55, 0.22, "Real-world trajectory frames provide the visual evidence for the hierarchy below.", 7.5, "muted", False, PP_ALIGN.CENTER)

    add_line(slide, 0.61, 2.10, 0.61, 4.45, color="line_dark", width=0.8, arrow=True)
    add_text(slide, 0.30, 2.08, 0.68, 0.22, "macro", 6.5, "muted", True, PP_ALIGN.CENTER)
    add_text(slide, 0.26, 4.47, 0.76, 0.22, "micro", 6.5, "muted", True, PP_ALIGN.CENTER)

    add_node(slide, 2.42, 2.04, 1.82, 0.33, "Goal-conditioned episode", "green_soft", "green", 7.4, True)
    add_node(slide, 1.33, 2.83, 1.54, 0.34, "Pick banana", "blue_soft", "blue", 7.4, True)
    add_node(slide, 3.82, 2.83, 1.54, 0.34, "Place in bowl", "amber_soft", "amber", 7.4, True)
    add_line(slide, 3.33, 2.37, 2.10, 2.83, color="line_dark", width=0.8, arrow=True)
    add_line(slide, 3.33, 2.37, 4.59, 2.83, color="line_dark", width=0.8, arrow=True)

    leaves = [
        (0.96, 3.76, "Approach", "blue_soft", "blue"),
        (1.82, 3.76, "Grasp", "blue_soft", "blue"),
        (2.68, 3.76, "Lift", "blue_soft", "blue"),
        (3.56, 3.76, "Move", "amber_soft", "amber"),
        (4.42, 3.76, "Align", "amber_soft", "amber"),
        (5.28, 3.76, "Release", "amber_soft", "amber"),
    ]
    for x, y, label, fill, line in leaves:
        add_node(slide, x, y, 0.70, 0.28, label, fill, line, 6.2, False)
    for x in [1.31, 2.17, 3.03]:
        add_line(slide, 2.10, 3.17, x, 3.76, color="line_dark", width=0.65, arrow=True)
    for x in [3.91, 4.77, 5.63]:
        add_line(slide, 4.59, 3.17, x, 3.76, color="line_dark", width=0.65, arrow=True)

    add_rect(slide, 1.14, 4.38, 4.84, 0.35, fill="white", line="line", radius=True)
    add_text(slide, 1.20, 4.43, 4.72, 0.21, "Long-horizon behavior naturally decomposes into reusable low-level controls.", 7.2, "ink", False, PP_ALIGN.CENTER)


def build_right_panel(slide):
    add_rect(slide, 6.54, 0.18, 6.62, 4.84, fill="white", line="line")
    add_text(slide, 6.68, 0.23, 6.24, 0.24, "(b) Memory topology: flat queue vs. hyperbolic space", 9.5, "ink", True)

    add_rect(slide, 6.76, 0.62, 6.16, 1.65, fill="panel_warm", line="line")
    add_rect(slide, 6.96, 0.78, 0.72, 0.22, fill="red_soft", line="red", radius=True)
    add_text(slide, 6.98, 0.805, 0.68, 0.13, "Baseline", 6.2, "red", True, PP_ALIGN.CENTER)
    add_text(slide, 7.77, 0.75, 2.05, 0.26, "Linear flat storage", 8.4, "ink", True)

    qx, qy = 7.00, 1.24
    for i in range(7):
        fill = "gray_soft" if i not in [1, 5] else "red_soft"
        add_rect(slide, qx + i * 0.52, qy, 0.36, 0.28, fill=fill, line="line", radius=True)
        label = f"f{i + 1}" if i < 5 else ("..." if i == 5 else "fn")
        add_text(slide, qx + i * 0.52, qy + 0.03, 0.36, 0.13, label, 6.1, "ink", False, PP_ALIGN.CENTER)
        if i < 6:
            add_line(slide, qx + i * 0.52 + 0.36, qy + 0.14, qx + (i + 1) * 0.52, qy + 0.14, color="line_dark", width=0.5)

    add_line(slide, 7.16, 1.76, 10.34, 1.76, color="red", width=0.9, arrow=True, dash=True)
    add_text(slide, 7.46, 1.84, 1.62, 0.18, "global scan", 6.5, "red", True, PP_ALIGN.CENTER)
    add_text(slide, 9.15, 1.84, 0.46, 0.18, "O(N)", 6.5, "red", True, PP_ALIGN.CENTER)
    add_cross(slide, 10.42, 1.70, 0.15, "red")
    add_text(slide, 10.64, 1.64, 1.72, 0.42, "hierarchy lost", 7.1, "red", True, PP_ALIGN.LEFT)

    add_rect(slide, 6.76, 2.52, 6.16, 2.25, fill="panel_green", line="line")
    add_rect(slide, 6.96, 2.68, 0.52, 0.22, fill="green_soft", line="green", radius=True)
    add_text(slide, 6.98, 2.705, 0.48, 0.13, "Ours", 6.2, "green", True, PP_ALIGN.CENTER)
    add_text(slide, 7.57, 2.65, 2.62, 0.26, "ECHO hyperbolic continuous space", 8.4, "ink", True)

    hyperbolic_asset = HYPERBOLIC_IMAGE2_CIRCLE if HYPERBOLIC_IMAGE2_CIRCLE.exists() else HYPERBOLIC_IMAGE2
    if hyperbolic_asset.exists():
        slide.shapes.add_picture(str(hyperbolic_asset), inch(7.20), inch(2.88), width=inch(1.86), height=inch(1.86))
    else:
        add_poincare_disk(slide, 8.13, 3.75, 0.78)
    add_text(slide, 7.46, 4.63, 1.36, 0.18, "Poincare disk", 6.4, "muted", False, PP_ALIGN.CENTER)

    add_node(slide, 9.38, 3.05, 1.18, 0.28, "semantic tree", "green_soft", "green", 6.4, True)
    add_node(slide, 10.72, 3.05, 1.42, 0.28, "top-down retrieval", "blue_soft", "blue", 6.4, True)
    add_line(slide, 9.00, 3.42, 9.38, 3.20, color="line_dark", width=0.75, arrow=True)
    add_line(slide, 10.56, 3.19, 10.72, 3.19, color="line_dark", width=0.75, arrow=True)

    add_rect(slide, 9.32, 3.66, 2.94, 0.48, fill="white", line="line", radius=True)
    add_text(slide, 9.43, 3.71, 2.72, 0.18, "Hyperbolic entailment cones", 7.0, "ink", True, PP_ALIGN.CENTER)
    add_text(slide, 9.43, 3.91, 2.72, 0.16, "encode partial and probabilistic containment", 6.2, "muted", False, PP_ALIGN.CENTER)
    add_line(slide, 8.98, 3.86, 9.32, 3.91, color="line_dark", width=0.75, arrow=True)

    add_rect(slide, 9.50, 4.34, 2.48, 0.26, fill="green_soft", line="green", radius=True)
    add_text(slide, 9.54, 4.37, 2.40, 0.14, "overlapping actions remain structured", 6.4, "green", True, PP_ALIGN.CENTER)


def build_presentation(frames: list[Path]) -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    safe_unlink(OUT_PPTX)
    prs = Presentation()
    prs.slide_width = inch(SLIDE_W)
    prs.slide_height = inch(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill="white", line=None)
    build_left_panel(slide, frames)
    build_right_panel(slide)
    prs.save(OUT_PPTX)


def export_with_powerpoint() -> None:
    try:
        import win32com.client
    except Exception as exc:
        print(f"PowerPoint export skipped: {exc}")
        return

    safe_unlink(OUT_PDF)
    safe_unlink(OUT_PNG)
    clear_export_dir(EXPORT_DIR)

    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = True
    app.DisplayAlerts = 0
    presentation = app.Presentations.Open(str(OUT_PPTX), False, True, False)
    try:
        presentation.SaveAs(str(OUT_PDF), 32)
        presentation.Export(str(EXPORT_DIR), "PNG", 3600, 1404)
    finally:
        presentation.Close()
        app.Quit()

    exported_pngs = sorted(EXPORT_DIR.glob("*.PNG")) + sorted(EXPORT_DIR.glob("*.png"))
    if exported_pngs:
        shutil.copy2(exported_pngs[0], OUT_PNG)


def verify_english_text() -> None:
    labels = [
        "(a) Natural hierarchy in embodied manipulation",
        "Real-world trajectory frames provide the visual evidence for the hierarchy below.",
        "Goal-conditioned episode",
        "Pick banana",
        "Place in bowl",
        "Approach",
        "Grasp",
        "Lift",
        "Move",
        "Align",
        "Release",
        "Long-horizon behavior naturally decomposes into reusable low-level controls.",
        "(b) Memory topology: flat queue vs. hyperbolic space",
        "Baseline",
        "Linear flat storage",
        "global scan",
        "hierarchy lost",
        "Ours",
        "ECHO hyperbolic continuous space",
        "Poincare disk",
        "semantic tree",
        "top-down retrieval",
        "Hyperbolic entailment cones",
        "encode partial and probabilistic containment",
        "overlapping actions remain structured",
    ]
    non_ascii = [label for label in labels if not label.isascii()]
    if non_ascii:
        raise RuntimeError(f"Non-ASCII annotation detected: {non_ascii}")
    joined = "\n".join(labels)
    if "Make Fruit Salad" in joined or "fruit salad" in joined.lower():
        raise RuntimeError("Forbidden fruit-salad title text detected.")


def main() -> None:
    verify_english_text()
    frames = extract_banana_frames()
    prepare_hyperbolic_image2()
    build_presentation(frames)
    export_with_powerpoint()
    print(f"PPTX: {OUT_PPTX}")
    print(f"PDF:  {OUT_PDF if OUT_PDF.exists() else 'not exported'}")
    print(f"PNG:  {OUT_PNG if OUT_PNG.exists() else 'not exported'}")


if __name__ == "__main__":
    main()
